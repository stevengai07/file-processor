# -*- coding: utf-8 -*-
"""Speaking-order Excel matching and ordered ZIP packaging."""

from __future__ import annotations

import io
import json
import os
import re
import uuid
import zipfile
from dataclasses import dataclass
from datetime import datetime
from difflib import SequenceMatcher
from typing import Any, Dict, Iterable, List, Optional, Tuple
from unicodedata import normalize as unicode_normalize

from openpyxl import Workbook, load_workbook

from schema import DocumentCandidate, MatchStatus, OrderItem, OrderMatch, OrderingTask


ORDER_ALIASES = {"顺序号", "顺序", "序号", "发言顺序", "排序", "order"}
TIME_ALIASES = {"时间", "发言时间", "汇报时间", "报告时间"}
TITLE_ALIASES = {"标题", "发言标题", "材料名称", "议题名称", "主题", "研发载体名称", "载体名称", "单位名称", "发言单位", "机构名称"}
REGION_ALIASES = {"所在地", "地区", "省市", "省份", "地域"}
CODE_ALIASES = {"编号", "材料编号", "议题编号", "文号", "文件编号"}
REPORT_HEADERS = [
    "顺序号", "时间", "顺序标题", "地区", "顺序编号", "原文件名", "提取标题", "提取编号",
    "新文件名", "匹配方式", "匹配分数", "状态", "原文依据", "是否人工修改",
]
_WEAK_TITLE_SUFFIXES = ("材料", "发言稿", "报告材料", "汇报材料", "汇报稿")
_ILLEGAL_FILENAME_CHARS = r'\\/:*?"<>|'


@dataclass
class ExcelParseResult:
    sheets: List[str]
    selected_sheet: str
    header_row: int
    column_mapping: Dict[str, Optional[str]]
    items: List[OrderItem]
    errors: List[str]
    rows: List[Dict[str, Any]]
    skipped_rows: int = 0
    auto_header_row: Optional[int] = None


def parse_order_excel(
    raw: bytes,
    filename: str = "order.xlsx",
    sheet_name: Optional[str] = None,
    header_row: int = 1,
    column_mapping: Optional[Dict[str, str]] = None,
) -> ExcelParseResult:
    if not filename.lower().endswith(".xlsx"):
        raise ValueError("顺序文件只支持 .xlsx")
    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sheets = wb.sheetnames
    selected_sheet = sheet_name or sheets[0]
    if selected_sheet not in sheets:
        raise ValueError(f"工作表不存在：{selected_sheet}")
    ws = wb[selected_sheet]
    auto_header_row = None
    if column_mapping is None:
        auto_header_row = _find_header_row(ws)
        if auto_header_row:
            header_row = auto_header_row
    headers = [str(c.value).strip() if c.value is not None else "" for c in ws[header_row]]
    mapping = _detect_mapping(headers)
    agenda_without_header = False
    if not column_mapping and not mapping.get("order_no") and _looks_like_agenda_without_header(ws):
        agenda_without_header = True
        mapping = {"order_no": "__col_1__", "time_range": "__col_2__", "title": "__col_3__", "region": "__col_4__", "code": None}
    if column_mapping:
        mapping.update({k: v for k, v in column_mapping.items() if v})

    errors: List[str] = []
    if not mapping.get("order_no"):
        errors.append("未识别顺序号列，请手动指定。")

    items: List[OrderItem] = []
    rows: List[Dict[str, Any]] = []
    seen_orders: Dict[int, int] = {}
    skipped_rows = 0

    index_by_header = {h: i for i, h in enumerate(headers) if h}
    min_row = 1 if agenda_without_header else header_row + 1
    for row_index, row in enumerate(ws.iter_rows(min_row=min_row, values_only=True), start=min_row):
        if not any(v not in (None, "") for v in row):
            continue

        order_raw = _value_from_row(row, mapping.get("order_no"), index_by_header)
        title = _clean_cell(_value_from_row(row, mapping.get("title"), index_by_header))
        code = _clean_cell(_value_from_row(row, mapping.get("code"), index_by_header))
        time_range = _clean_cell(_value_from_row(row, mapping.get("time_range"), index_by_header))
        region = _clean_cell(_value_from_row(row, mapping.get("region"), index_by_header))
        row_view = {"source_row": row_index, "order_no": order_raw, "time_range": time_range, "title": title, "region": region, "code": code, "valid": True, "error": ""}

        order_no = _parse_order_no(order_raw)
        if order_no is None:
            skipped_rows += 1
            continue
        elif not title and not code:
            row_view.update(valid=False, error="标题和编号不能同时为空")
            errors.append(f"第 {row_index} 行：标题和编号不能同时为空。")
        else:
            if order_no in seen_orders:
                row_view.update(valid=False, error=f"顺序号重复，首次出现于第 {seen_orders[order_no]} 行")
                errors.append(f"第 {row_index} 行：顺序号 {order_no} 重复，首次出现于第 {seen_orders[order_no]} 行。")
            else:
                seen_orders[order_no] = row_index
                items.append(OrderItem(
                    item_id=str(uuid.uuid4()),
                    order_no=order_no,
                    title=title or None,
                    code=code or None,
                    source_row=row_index,
                    time_range=time_range or None,
                    region=region or None,
                ))
        rows.append(row_view)

    items.sort(key=lambda x: x.order_no)
    return ExcelParseResult(sheets, selected_sheet, header_row, mapping, items, errors, rows, skipped_rows, auto_header_row)


def normalise_code(value: Optional[str]) -> str:
    if not value:
        return ""
    text = unicode_normalize("NFKC", str(value)).upper()
    return "".join(ch for ch in text if ch.isalnum())


def normalise_title(value: Optional[str]) -> str:
    if not value:
        return ""
    text = unicode_normalize("NFKC", str(value))
    text = re.sub(r"\.(pdf|docx?|xlsx)$", "", text, flags=re.I)
    text = re.sub(r"[《》<>“”\"'‘’（）()\[\]【】]", "", text)
    text = re.sub(r"\s+", "", text).strip(" .。")
    for suffix in _WEAK_TITLE_SUFFIXES:
        if text.endswith(suffix):
            text = text[: -len(suffix)]
    return text


def calculate_match_score(order: OrderItem, doc: DocumentCandidate) -> Tuple[float, str, str]:
    order_code = order.code or ""
    doc_code = doc.extracted_code or ""
    if order_code and doc_code:
        if order_code == doc_code:
            return 1.0, "code_exact", "编号完全一致"
        if normalise_code(order_code) and normalise_code(order_code) == normalise_code(doc_code):
            return 0.98, "code_normalised", "编号标准化后一致"

    order_title = order.title or ""
    doc_title = doc.extracted_title or _strip_extension(doc.filename)
    if order_title and doc_title:
        if order_title == doc_title:
            return 0.95, "title_exact", "标题完全一致"
        ot = normalise_title(order_title)
        dt = normalise_title(doc_title)
        if ot and dt:
            if ot in dt or dt in ot:
                return 0.90, "title_contains", "标题包含匹配"
            sim = SequenceMatcher(None, ot, dt).ratio()
            if sim >= 0.90:
                return 0.88, "title_similarity", f"标题相似度 {sim:.2f}"
            if sim >= 0.80:
                return 0.75, "title_similarity", f"标题相似度 {sim:.2f}"
            keyword_score = _keyword_score(ot, dt)
            if keyword_score:
                return min(keyword_score, 0.60), "keyword", "少量关键词命中"
    return 0.0, "none", "无候选匹配依据"


def assign_unique_matches(items: List[OrderItem], documents: List[DocumentCandidate]) -> List[OrderMatch]:
    used: set[str] = set()
    matches: List[OrderMatch] = []
    for item in sorted(items, key=lambda x: x.order_no):
        scored = []
        for doc in documents:
            score, method, reason = calculate_match_score(item, doc)
            if score > 0:
                scored.append((score, method, reason, doc.file_id))
        scored.sort(key=lambda x: x[0], reverse=True)
        match = resolve_match(item, scored, used)
        if match.selected_file_id:
            used.add(match.selected_file_id)
        matches.append(match)
    return matches


def resolve_match(item: OrderItem, scored_candidates: List[Tuple[float, str, str, str]], used_file_ids: set[str]) -> OrderMatch:
    candidate_ids = [c[3] for c in scored_candidates[:5]]
    if not scored_candidates:
        return OrderMatch(
            match_id=str(uuid.uuid4()), order_item_id=item.item_id,
            candidate_file_ids=[], status=MatchStatus.UNMATCHED, match_reason="没有候选文件",
        )
    top = scored_candidates[0]
    second_score = scored_candidates[1][0] if len(scored_candidates) > 1 else 0.0
    occupied = top[3] in used_file_ids
    can_auto = top[0] >= 0.85 and (top[0] - second_score) >= 0.10 and not occupied
    return OrderMatch(
        match_id=str(uuid.uuid4()),
        order_item_id=item.item_id,
        selected_file_id=top[3] if can_auto else None,
        candidate_file_ids=candidate_ids,
        score=round(top[0], 2),
        status=MatchStatus.MATCHED if can_auto else MatchStatus.NEEDS_REVIEW,
        match_method=top[1],
        match_reason=top[2] if not occupied else "最高候选已被其他顺序项占用",
        manually_confirmed=False,
    )


def confirm_match(match: OrderMatch, file_id: Optional[str], skip: bool = False) -> OrderMatch:
    if skip:
        match.selected_file_id = None
        match.status = MatchStatus.SKIPPED
        match.manually_confirmed = True
        match.match_reason = "人工跳过"
        return match
    if not file_id:
        raise ValueError("请选择文件后再确认")
    match.selected_file_id = file_id
    match.status = MatchStatus.CONFIRMED
    match.manually_confirmed = True
    match.match_method = "manual"
    match.match_reason = "人工确认"
    match.score = max(match.score, 1.0)
    return match


def validate_assignments(matches: List[OrderMatch]) -> List[str]:
    errors: List[str] = []
    seen: Dict[str, str] = {}
    for match in matches:
        if match.status in (MatchStatus.NEEDS_REVIEW, MatchStatus.UNMATCHED, MatchStatus.FAILED):
            errors.append(f"匹配 {match.match_id[:8]} 尚未处理：{match.status.value}")
        if match.selected_file_id:
            if match.selected_file_id in seen:
                errors.append(f"文件重复分配：{match.selected_file_id}")
            seen[match.selected_file_id] = match.match_id
    return errors


def extract_document_identity(
    filename: str,
    raw: bytes,
    file_id: Optional[str] = None,
    ocr_preset: str = "scanner",
    ocr_lang: str = "chi_sim+eng",
    use_ollama: bool = True,
) -> DocumentCandidate:
    ext = os.path.splitext(filename)[1].lower()
    candidate = DocumentCandidate(
        file_id=file_id or str(uuid.uuid4()),
        filename=filename,
        extension=ext,
        size=len(raw),
    )
    title_hint, code_hint = _identity_from_filename(filename)
    candidate.extracted_title = title_hint
    candidate.extracted_code = code_hint
    try:
        text = ""
        if ext == ".pptx":
            text = _extract_pptx_text(raw)
        elif ext == ".ppt":
            text = ""
            candidate.extraction_error = "旧版 .ppt 暂不支持正文提取，已使用文件名参与匹配。"
        else:
            from extractor import extract
            doc = extract(filename, raw, {
                "ocr_enabled": ocr_preset not in {"off", "关闭"},
                "ocr_preset": ocr_preset,
                "tesseract_lang": ocr_lang,
            })
            text = doc.plain_text or doc.full_text
        if use_ollama and text.strip():
            data = _extract_identity_with_ollama(text[:12000])
            candidate.extracted_title = data.get("document_title") or candidate.extracted_title
            candidate.extracted_code = data.get("document_code") or candidate.extracted_code
            candidate.speaker_or_department = data.get("speaker_or_department")
            candidate.source_snippet = data.get("source_snippet")
    except Exception as exc:
        candidate.extraction_error = str(exc)
    return candidate


def _extract_pptx_text(raw: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx 未安装，无法提取 PPTX 正文")

    prs = Presentation(io.BytesIO(raw))
    parts: List[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_text.append(" | ".join(cells))
        if slide_text:
            parts.append(f"--- Slide {slide_index} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def build_matching_report(task: OrderingTask) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "匹配清单"
    ws.append(REPORT_HEADERS)
    docs = {d.file_id: d for d in task.documents}
    items = {i.item_id: i for i in task.items}
    used_names: set[str] = {m.output_filename for m in task.matches if m.output_filename}
    for match in sorted(task.matches, key=lambda m: items[m.order_item_id].order_no if m.order_item_id in items else 999999):
        item = items.get(match.order_item_id)
        doc = docs.get(match.selected_file_id or "")
        if item and doc and not match.output_filename and match.status != MatchStatus.SKIPPED:
            match.output_filename = make_output_filename(item, doc.extension, used_names)
        ws.append([
            item.order_no if item else "",
            item.time_range if item else "",
            item.title if item else "",
            item.region if item else "",
            item.code if item else "",
            doc.filename if doc else "",
            doc.extracted_title if doc else "",
            doc.extracted_code if doc else "",
            match.output_filename or "",
            match.match_method,
            match.score,
            match.status.value,
            doc.source_snippet if doc else "",
            "是" if match.manually_confirmed else "否",
        ])
    for col in ws.columns:
        letter = col[0].column_letter
        width = min(max(len(str(cell.value or "")) for cell in col) + 3, 50)
        ws.column_dimensions[letter].width = width
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_renamed_zip(task: OrderingTask, file_bytes: Dict[str, bytes]) -> bytes:
    errors = validate_assignments(task.matches)
    if errors:
        raise ValueError("导出前请先处理问题：" + "；".join(errors))
    docs = {d.file_id: d for d in task.documents}
    items = {i.item_id: i for i in task.items}
    used_names: set[str] = set()
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for match in sorted(task.matches, key=lambda m: items[m.order_item_id].order_no if m.order_item_id in items else 999999):
            if match.status == MatchStatus.SKIPPED or not match.selected_file_id:
                continue
            item = items[match.order_item_id]
            doc = docs[match.selected_file_id]
            raw = file_bytes.get(doc.file_id)
            if raw is None:
                raise ValueError(f"缺少源文件内容：{doc.filename}")
            output_name = make_output_filename(item, doc.extension, used_names)
            match.output_filename = output_name
            zf.writestr(output_name, raw)
        zf.writestr("匹配清单.xlsx", build_matching_report(task))
    return buf.getvalue()


def make_output_filename(item: OrderItem, extension: str, used_names: Optional[set[str]] = None) -> str:
    used_names = used_names if used_names is not None else set()
    title = item.title or item.code or f"材料{item.order_no}"
    safe_title = sanitize_filename(title, max_len=120)
    ext = extension if extension.startswith(".") else f".{extension}"
    base = f"{item.order_no:03d}_{safe_title}"
    name = f"{base}{ext}"
    counter = 2
    while name in used_names:
        name = f"{base}_{counter}{ext}"
        counter += 1
    used_names.add(name)
    return name


def sanitize_filename(value: str, max_len: int = 120) -> str:
    text = unicode_normalize("NFKC", str(value))
    text = re.sub(f"[{re.escape(_ILLEGAL_FILENAME_CHARS)}]", "_", text)
    text = re.sub(r"\s+", " ", text).strip(" .")
    text = text.replace("..", ".")
    if not text:
        text = "未命名"
    return text[:max_len].strip(" .") or "未命名"


def create_ordering_task(name: str, order_filename: str, items: List[OrderItem]) -> OrderingTask:
    return OrderingTask(ordering_id=str(uuid.uuid4()), name=name, order_filename=order_filename, items=items)


def _detect_mapping(headers: List[str]) -> Dict[str, Optional[str]]:
    lowered = {h.strip().lower(): h for h in headers if h}
    return {
        "order_no": _find_alias(lowered, ORDER_ALIASES),
        "time_range": _find_alias(lowered, TIME_ALIASES),
        "title": _find_alias(lowered, TITLE_ALIASES),
        "region": _find_alias(lowered, REGION_ALIASES),
        "code": _find_alias(lowered, CODE_ALIASES),
    }


def _find_alias(lowered: Dict[str, str], aliases: Iterable[str]) -> Optional[str]:
    for alias in aliases:
        if alias.lower() in lowered:
            return lowered[alias.lower()]
    for key, original in lowered.items():
        if any(alias.lower() in key for alias in aliases):
            return original
    return None


def _value_from_row(row: Tuple[Any, ...], header: Optional[str], index_by_header: Dict[str, int]) -> Any:
    if not header or header not in index_by_header:
        if header and header.startswith("__col_") and header.endswith("__"):
            try:
                idx = int(header.removeprefix("__col_").removesuffix("__")) - 1
            except ValueError:
                return None
            return row[idx] if idx < len(row) else None
        return None
    idx = index_by_header[header]
    return row[idx] if idx < len(row) else None


def _clean_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _parse_order_no(value: Any) -> Optional[int]:
    if value in (None, ""):
        return None
    try:
        return int(float(str(value).strip()))
    except ValueError:
        return None


def _looks_like_agenda_without_header(ws) -> bool:
    checked = 0
    matched = 0
    for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 5), values_only=True):
        if not any(v not in (None, "") for v in row):
            continue
        checked += 1
        if len(row) >= 3 and _parse_order_no(row[0]) is not None and _clean_cell(row[2]):
            matched += 1
    return checked > 0 and matched == checked


def _find_header_row(ws, max_scan_rows: int = 30) -> Optional[int]:
    best_row = None
    best_score = 0
    max_row = min(ws.max_row, max_scan_rows)
    for row_index, row in enumerate(ws.iter_rows(min_row=1, max_row=max_row, values_only=True), start=1):
        values = [_clean_cell(v) for v in row]
        if not any(values):
            continue
        mapping = _detect_mapping(values)
        score = 0
        if mapping.get("order_no"):
            score += 3
        if mapping.get("title"):
            score += 3
        if mapping.get("time_range"):
            score += 1
        if mapping.get("region"):
            score += 1
        if mapping.get("code"):
            score += 1
        if score > best_score:
            best_score = score
            best_row = row_index
    return best_row if best_score >= 6 else None


def _strip_extension(filename: str) -> str:
    return os.path.splitext(os.path.basename(filename))[0]


def _identity_from_filename(filename: str) -> Tuple[str, Optional[str]]:
    title = _strip_extension(filename)
    code_match = re.search(r"([A-Za-z]{0,6}\d[\w\-〔〕\[\]（）()号第]*\d*)", unicode_normalize("NFKC", title))
    code = code_match.group(1) if code_match else None
    clean_title = re.sub(r"^[\d一二三四五六七八九十]+[、._\-\s]+", "", title).strip()
    return clean_title, code


def _keyword_score(left: str, right: str) -> float:
    tokens = [t for t in re.split(r"[，,。；;：:\s]+", left) if len(t) >= 2]
    if not tokens:
        return 0.0
    hits = sum(1 for t in tokens if t in right)
    return 0.60 * hits / len(tokens) if hits else 0.0


def _extract_identity_with_ollama(text: str) -> Dict[str, Any]:
    from langchain_core.messages import HumanMessage, SystemMessage
    from langchain_openai import ChatOpenAI

    llm = ChatOpenAI(
        model="qwen3.6:latest",
        base_url=os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434/v1"),
        api_key="ollama",
        temperature=0,
        max_tokens=1024,
        timeout=300,
        max_retries=1,
    )
    system = "只从原文提取文档标题、材料编号、发言单位。不得猜测。找不到返回 null。只返回 JSON。"
    prompt = f"""返回 JSON：{{"document_title": null, "document_code": null, "speaker_or_department": null, "source_snippet": null}}

原文：
{text}
"""
    response = llm.invoke([SystemMessage(content=system), HumanMessage(content=prompt)])
    raw = str(response.content).strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw).strip()
    data = json.loads(raw)
    if not isinstance(data, dict):
        raise ValueError("Ollama 返回不是 JSON 对象")
    return data
