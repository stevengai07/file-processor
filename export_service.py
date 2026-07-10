#!/usr/bin/env python3
"""
export_service.py — Unified Export Engine
Contains builders for Excel, Word (DOCX), and PowerPoint (PPTX).
"""

from __future__ import annotations
import io
import re
import datetime
from typing import Any, List, Dict

# ──────────────────────────────────────────────────────────────────────────────
# 1. EXCEL EXPORT (For Page 3: Batch Extraction Results)
# ──────────────────────────────────────────────────────────────────────────────
def build_excel(task: Any, template: Any, results_to_export: List[Any], include_log: bool = True) -> bytes:
    import pandas as pd
    
    data = []
    field_keys = [f.key for f in template.fields]
    field_names = [f.name for f in template.fields]

    for r in results_to_export:
        status_val = r.status.value if hasattr(r, 'status') and hasattr(r.status, 'value') else str(getattr(r, 'status', ''))
        
        row = {
            "文件名": getattr(r, "filename", ""),
            "状态": status_val,
            "耗时(s)": round(getattr(r, "elapsed_seconds", 0.0) or 0.0, 1),
        }
        
        f_dict = {fv.key: fv for fv in getattr(r, "fields", [])}
        for f_key, f_name in zip(field_keys, field_names):
            if f_key in f_dict:
                val = f_dict[f_key].value
                if isinstance(val, list):
                    val = ", ".join(map(str, val))
                elif isinstance(val, bool):
                    val = "是" if val else "否"
                row[f_name] = str(val) if val is not None else ""
            else:
                row[f_name] = ""
                
        if include_log:
            issues = getattr(r, "issues", [])
            row["日志警告"] = "\n".join([f"[{i.field_name}] {i.message}" for i in issues]) if issues else "无异常"
            
        data.append(row)

    df = pd.DataFrame(data)
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='数据提取明细')
        worksheet = writer.sheets['数据提取明细']
        for idx, col in enumerate(df.columns):
            max_len = max((df[col].astype(str).map(len).max(), len(col))) + 2
            worksheet.column_dimensions[chr(65 + idx)].width = min(max_len, 50)
            
    return buf.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# 2. DOCX EXPORT (For Page 3: Batch Extraction Results)
# ──────────────────────────────────────────────────────────────────────────────
def build_docx(task: Any, template: Any, results_to_export: List[Any], include_log: bool = True) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    
    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'
    try:
        from docx.oxml.ns import qn
        style.font.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
    except Exception:
        pass
    style.font.size = Pt(10.5)

    title = doc.add_heading('批量文档智能萃取报告', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    task_name = getattr(task, "name", "未知任务")
    p_meta = doc.add_paragraph()
    p_meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_meta.add_run(f"任务名称: {task_name} | 导出时间: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}")
    run.font.color.rgb = RGBColor(120, 120, 120)
    run.font.size = Pt(9)
    
    for idx, r in enumerate(results_to_export):
        h1 = doc.add_heading(f"{idx+1}. {getattr(r, 'filename', 'Unknown')}", level=1)
        h1.runs[0].font.color.rgb = RGBColor(75, 46, 131)
        
        status_val = r.status.value if hasattr(r, 'status') and hasattr(r.status, 'value') else str(getattr(r, 'status', ''))
        elapsed = getattr(r, "elapsed_seconds", 0)
        doc.add_paragraph(f"解析状态: {status_val}   耗时: {round(elapsed or 0, 1)}s")

        fields_list = getattr(r, "fields", [])
        if fields_list:
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = '字段名称'
            hdr_cells[1].text = '提取内容'
            
            for fv in fields_list:
                row_cells = table.add_row().cells
                row_cells[0].text = str(fv.name)
                val = fv.value
                if isinstance(val, list):
                    val = ", ".join(map(str, val))
                elif isinstance(val, bool):
                    val = "是" if val else "否"
                row_cells[1].text = str(val) if val is not None else ""
        
        if include_log:
            issues = getattr(r, "issues", [])
            if issues:
                doc.add_heading("⚠️ 预警与问题", level=3)
                for issue in issues:
                    doc.add_paragraph(f"[{issue.field_name}] {issue.message}", style='List Bullet')
        
        if idx < len(results_to_export) - 1:
            doc.add_page_break()

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
# 3. PPTX EXPORT (For Page 4: AI Console Presentation)
# ──────────────────────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    pass

# Design tokens
W = Inches(13.33)
H = Inches(7.5)

C_PURPLE_DARK  = RGBColor(0x2D, 0x17, 0x6B)
C_PURPLE_MID   = RGBColor(0x4A, 0x2C, 0x99)
C_PURPLE_LIGHT = RGBColor(0x7C, 0x5C, 0xC8)
C_TEAL         = RGBColor(0x00, 0x96, 0x88)
C_AMBER        = RGBColor(0xF5, 0x9E, 0x0E)
C_GREEN        = RGBColor(0x16, 0xA3, 0x4A)
C_RED          = RGBColor(0xDC, 0x26, 0x26)
C_GRAY_TEXT    = RGBColor(0x37, 0x41, 0x51)
C_MUTED        = RGBColor(0x6B, 0x72, 0x80)
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT     = RGBColor(0xF8, 0xF7, 0xFF)
C_CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)
C_DIVIDER      = RGBColor(0xE5, 0xE7, 0xEB)

FONT_DISPLAY = "Calibri"
FONT_BODY    = "Calibri"

BULLETS_PER_SLIDE  = 7
CHARS_PER_SLIDE    = 900
PILLS_PER_SLIDE    = 12
ITEMS_PER_LIST_SLIDE = 10
ACTIONS_PER_SLIDE  = 7

def _strip_md(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*",       r"\1", text)
    text = re.sub(r"`(.+?)`",           r"\1", text)
    text = re.sub(r"^#{1,6}\s+",       "",  text, flags=re.M)
    text = re.sub(r"^[-–—]{3,}$",       "",  text, flags=re.M)
    text = re.sub(r"^[\*\-]\s+",     "",  text, flags=re.M)
    text = re.sub(r"\n{3,}",           "\n\n", text)
    return text.strip()

def _chunk_list(lst: list, size: int) -> list:
    return [lst[i:i+size] for i in range(0, max(len(lst), 1), size)]

def _chunk_text(text: str, max_chars: int) -> list:
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks, current, current_len = [], [], 0
    for para in paragraphs:
        if current_len + len(para) > max_chars and current:
            chunks.append("\n".join(current))
            current, current_len = [], 0
        current.append(para)
        current_len += len(para)
    if current:
        chunks.append("\n".join(current))
    return chunks or [text]

def _slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def _text_box(slide, x, y, w, h, text: str, *,
              font_name=FONT_BODY, font_size=14, bold=False, italic=False,
              color: RGBColor = C_GRAY_TEXT, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _strip_md(text)
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def _text_box_multiline(slide, x, y, w, h, lines: list, *,
                         font_name=FONT_BODY, font_size=13,
                         color=C_GRAY_TEXT, bullet_color=C_PURPLE_MID,
                         indent=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if indent:
            from pptx.util import Pt as _Pt
            p.space_before = _Pt(3)
        run = p.add_run()
        run.text = ("• " if indent else "") + _strip_md(line)
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb

def _footer(slide, label: str = "AI Document Analysis"):
    bar_h = Inches(0.28)
    _rect(slide, 0, H - bar_h, W, bar_h, fill=C_PURPLE_DARK)
    _text_box(slide, Inches(0.3), H - bar_h, Inches(6), bar_h,
              label, font_size=8, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.LEFT)
    ts = datetime.datetime.now().strftime("%Y-%m-%d")
    _text_box(slide, W - Inches(2), H - bar_h, Inches(1.9), bar_h,
              ts, font_size=8, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.RIGHT)

def _section_heading(slide, title: str, y_pos, subtitle: str = ""):
    _rect(slide, Inches(0.4), y_pos, Inches(0.06), Inches(0.38), fill=C_PURPLE_MID)
    _text_box(slide, Inches(0.6), y_pos - Inches(0.02), Inches(11.5), Inches(0.45),
              title, font_name=FONT_DISPLAY, font_size=18, bold=True,
              color=C_PURPLE_DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        _text_box(slide, Inches(0.6), y_pos + Inches(0.42), Inches(11.5), Inches(0.32),
                  subtitle, font_size=10, color=C_MUTED)

def _content_card(slide, x, y, w, h, accent_color=None):
    _rect(slide, x, y, w, h, fill=C_CARD_BG, line_color=C_DIVIDER, line_width=1)
    if accent_color:
        _rect(slide, x, y, w, Inches(0.055), fill=accent_color)

def _slide_cover(prs, info: dict, source_filename: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, C_PURPLE_DARK)
    _rect(slide, 0, 0, W, Inches(0.08), fill=C_PURPLE_LIGHT)
    _rect(slide, 0, Inches(0.08), W, Inches(0.04), fill=C_AMBER)

    _text_box(slide, Inches(0.7), Inches(0.9), Inches(11.9), Inches(1.1),
              "AI 文档分析报告",
              font_name=FONT_DISPLAY, font_size=40, bold=True,
              color=C_WHITE, align=PP_ALIGN.LEFT)

    clean_src = source_filename.replace("_", " ")
    _text_box(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.55),
              f"来源文件：{clean_src}",
              font_size=13, color=RGBColor(0xCC, 0xBB, 0xFF))

    ts = datetime.datetime.now().strftime("导出时间：%Y-%m-%d  %H:%M")
    _text_box(slide, Inches(0.7), Inches(2.55), Inches(6), Inches(0.4),
              ts, font_size=11, color=RGBColor(0xAA, 0x99, 0xDD))

    stats = [
        ("📄", str(info.get("page_count", "—")), "页数"),
        ("🔤", f"{info.get('char_count', 0):,}", "字符数"),
        ("🏷", str(len(info.get("topics", []))), "主题"),
        ("✅", str(len(info.get("actions", []))), "行动项"),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.35)
    gap    = Inches(0.28)
    start_x = (W - (card_w * 4 + gap * 3)) / 2
    y = Inches(3.6)

    for i, (emoji, value, label) in enumerate(stats):
        cx = start_x + i * (card_w + gap)
        _rect(slide, cx, y, card_w, card_h,
              fill=RGBColor(0x3D, 0x27, 0x7B),
              line_color=C_PURPLE_LIGHT, line_width=1)
        _text_box(slide, cx, y + Inches(0.08), card_w, Inches(0.45),
                  emoji, font_size=22, align=PP_ALIGN.CENTER, color=C_WHITE)
        _text_box(slide, cx, y + Inches(0.48), card_w, Inches(0.52),
                  value, font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _text_box(slide, cx, y + Inches(0.98), card_w, Inches(0.32),
                  label, font_size=11, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.CENTER)

def _slides_summary(prs, info: dict):
    summary = _strip_md(info.get("summary", "暂无摘要。"))
    chunks = _chunk_text(summary, CHARS_PER_SLIDE)

    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _slide_bg(slide, C_BG_LIGHT)
        subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
        _section_heading(slide, "文档摘要", Inches(0.38), subtitle)

        card_x, card_y = Inches(0.4), Inches(1.0)
        card_w = W - Inches(0.8)
        card_h = H - Inches(1.7)
        _content_card(slide, card_x, card_y, card_w, card_h, C_PURPLE_MID)
        _text_box(slide, card_x + Inches(0.25), card_y + Inches(0.18),
                  card_w - Inches(0.5), card_h - Inches(0.3),
                  chunk, font_size=15, color=C_GRAY_TEXT, word_wrap=True)
        _footer(slide)

def _slides_topics(prs, info: dict):
    topics = info.get("topics", [])
    if not topics:
        return
    chunks = _chunk_list(topics, PILLS_PER_SLIDE)
    pill_colors = [C_PURPLE_MID, C_TEAL, C_AMBER,
                   RGBColor(0x06, 0x82, 0x72), RGBColor(0x15, 0x78, 0xC2)]
    pill_w, pill_h = Inches(3.8), Inches(0.68)
    cols = 3
    gap_x, gap_y = Inches(0.28), Inches(0.32)
    start_x = (W - (pill_w * cols + gap_x * (cols - 1))) / 2

    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _slide_bg(slide, C_BG_LIGHT)
        subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
        _section_heading(slide, "核心主题", Inches(0.38), subtitle)

        start_y = Inches(1.15)
        for i, topic in enumerate(chunk):
            row, col = divmod(i, cols)
            px = start_x + col * (pill_w + gap_x)
            py = start_y + row * (pill_h + gap_y)
            c = pill_colors[i % len(pill_colors)]
            _rect(slide, px, py, pill_w, pill_h, fill=c)
            _text_box(slide, px, py, pill_w, pill_h,
                      _strip_md(topic), font_size=14, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)
        _footer(slide)

def _slides_entities_dates(prs, info: dict):
    entities = info.get("entities", [])
    dates    = info.get("dates", [])
    if not entities and not dates:
        return

    e_chunks = _chunk_list(entities, ITEMS_PER_LIST_SLIDE) if entities else [[]]
    d_chunks = _chunk_list(dates,    ITEMS_PER_LIST_SLIDE) if dates    else [[]]
    n_slides = max(len(e_chunks), len(d_chunks))

    for idx in range(n_slides):
        e_chunk = e_chunks[idx] if idx < len(e_chunks) else []
        d_chunk = d_chunks[idx] if idx < len(d_chunks) else []

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _slide_bg(slide, C_BG_LIGHT)
        subtitle = f"({idx+1}/{n_slides})" if n_slides > 1 else ""
        _section_heading(slide, "实体与时间线", Inches(0.38), subtitle)

        half_w = W / 2 - Inches(0.6)
        card_h = H - Inches(1.7)
        card_y = Inches(1.0)

        if e_chunk:
            _content_card(slide, Inches(0.4), card_y, half_w, card_h, C_PURPLE_MID)
            _text_box(slide, Inches(0.5), card_y + Inches(0.1), half_w - Inches(0.2), Inches(0.42),
                      "📌 关键实体", font_size=13, bold=True, color=C_PURPLE_DARK)
            for i, ent in enumerate(e_chunk):
                ey = card_y + Inches(0.58) + i * Inches(0.5)
                if ey + Inches(0.42) > card_y + card_h - Inches(0.15):
                    break
                _rect(slide, Inches(0.5), ey, Inches(0.06), Inches(0.32), fill=C_PURPLE_MID)
                _text_box(slide, Inches(0.68), ey - Inches(0.03),
                          half_w - Inches(0.35), Inches(0.42),
                          _strip_md(ent), font_size=12, color=C_GRAY_TEXT)

        if d_chunk:
            dx = W / 2 + Inches(0.2)
            _content_card(slide, dx, card_y, half_w, card_h, C_TEAL)
            _text_box(slide, dx + Inches(0.1), card_y + Inches(0.1),
                      half_w - Inches(0.2), Inches(0.42),
                      "📅 时间节点", font_size=13, bold=True, color=C_TEAL)
            for i, dt in enumerate(d_chunk):
                dy_pos = card_y + Inches(0.58) + i * Inches(0.5)
                if dy_pos + Inches(0.42) > card_y + card_h - Inches(0.15):
                    break
                _rect(slide, dx + Inches(0.1), dy_pos, Inches(0.06), Inches(0.32), fill=C_TEAL)
                _text_box(slide, dx + Inches(0.28), dy_pos - Inches(0.03),
                          half_w - Inches(0.35), Inches(0.42),
                          _strip_md(dt), font_size=12, color=C_GRAY_TEXT)
        _footer(slide)

def _slides_actions(prs, info: dict):
    actions = info.get("actions", [])
    if not actions:
        return

    chunks = _chunk_list(actions, ACTIONS_PER_SLIDE)
    badge_colors = [C_PURPLE_MID, C_TEAL, C_AMBER, C_GREEN,
                    C_RED, C_PURPLE_LIGHT, C_TEAL]

    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _slide_bg(slide, C_BG_LIGHT)
        subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
        _section_heading(slide, "行动项与建议", Inches(0.38), subtitle)

        row_h   = Inches(0.68)
        row_w   = W - Inches(0.8)
        start_y = Inches(1.05)
        global_i = idx * ACTIONS_PER_SLIDE 

        for i, action in enumerate(chunk):
            ry = start_y + i * (row_h + Inches(0.1))
            if ry + row_h > H - Inches(0.5):
                break
            bc = badge_colors[(global_i + i) % len(badge_colors)]
            _content_card(slide, Inches(0.4), ry, row_w, row_h)
            _rect(slide, Inches(0.4), ry, Inches(0.07), row_h, fill=bc)
            _rect(slide, Inches(0.57), ry + Inches(0.17), Inches(0.36), Inches(0.36), fill=bc)
            _text_box(slide, Inches(0.57), ry + Inches(0.17), Inches(0.36), Inches(0.36),
                      str(global_i + i + 1), font_size=12, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)
            _text_box(slide, Inches(1.05), ry + Inches(0.12),
                      row_w - Inches(0.72), Inches(0.50),
                      _strip_md(action), font_size=13, color=C_GRAY_TEXT)
        _footer(slide)

def _slides_extra_sections(prs, info: dict):
    sections = info.get("sections", [])
    if not sections:
        return

    for section in sections:
        title   = _strip_md(section.get("title", "附加内容"))
        content = section.get("content", "")
        stype   = section.get("type", "text")

        if stype == "text":
            text = _strip_md(content) if isinstance(content, str) else ""
            chunks = _chunk_text(text, CHARS_PER_SLIDE)
            for idx, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _slide_bg(slide, C_BG_LIGHT)
                subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
                _section_heading(slide, title, Inches(0.38), subtitle)
                card_x, card_y = Inches(0.4), Inches(1.0)
                card_w = W - Inches(0.8)
                card_h = H - Inches(1.7)
                _content_card(slide, card_x, card_y, card_w, card_h, C_PURPLE_MID)
                _text_box(slide, card_x + Inches(0.25), card_y + Inches(0.18),
                          card_w - Inches(0.5), card_h - Inches(0.3),
                          chunk, font_size=14, color=C_GRAY_TEXT, word_wrap=True)
                _footer(slide)

        elif stype == "bullets":
            items = content if isinstance(content, list) else [
                l.strip() for l in str(content).splitlines() if l.strip()
            ]
            chunks = _chunk_list(items, BULLETS_PER_SLIDE)
            for idx, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _slide_bg(slide, C_BG_LIGHT)
                subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
                _section_heading(slide, title, Inches(0.38), subtitle)
                card_x, card_y = Inches(0.4), Inches(1.0)
                card_w = W - Inches(0.8)
                card_h = H - Inches(1.7)
                _content_card(slide, card_x, card_y, card_w, card_h, C_AMBER)
                _text_box_multiline(
                    slide,
                    card_x + Inches(0.3), card_y + Inches(0.15),
                    card_w - Inches(0.5), card_h - Inches(0.25),
                    chunk, font_size=14, color=C_GRAY_TEXT,
                    bullet_color=C_AMBER
                )
                _footer(slide)

        elif stype == "table":
            rows_data = content if isinstance(content, list) else []
            if not rows_data: continue
            
            num_cols = max(len(r) for r in rows_data) if rows_data else 1
            chunks = _chunk_list(rows_data, 8) 
            
            for idx, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _slide_bg(slide, C_BG_LIGHT)
                subtitle = f"({idx+1}/{len(chunks)})" if len(chunks) > 1 else ""
                _section_heading(slide, title, Inches(0.38), subtitle)

                table_x, table_y = Inches(0.5), Inches(1.1)
                table_w, table_h = W - Inches(1.0), Inches(0.5 * len(chunk))
                
                table_shape = slide.shapes.add_table(len(chunk), num_cols, table_x, table_y, table_w, table_h)
                table = table_shape.table
                
                for r_idx, row in enumerate(chunk):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = _strip_md(str(val))
                        
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                run.font.name = FONT_BODY
                                run.font.size = Pt(13)
                                if r_idx == 0 or c_idx == 0:
                                    run.font.bold = True
                                    run.font.color.rgb = C_PURPLE_DARK
                                else:
                                    run.font.color.rgb = C_GRAY_TEXT

                _footer(slide)

def _slide_closing(prs, source_filename: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, C_PURPLE_DARK)
    _rect(slide, W - Inches(2.8), Inches(-0.5), Inches(3.5), Inches(3.5),
          fill=RGBColor(0x3D, 0x27, 0x7B))
    _rect(slide, Inches(-0.8), H - Inches(2.2), Inches(3), Inches(3),
          fill=RGBColor(0x3D, 0x27, 0x7B))
    _text_box(slide, Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.3),
              "分析完成",
              font_name=FONT_DISPLAY, font_size=52, bold=True,
              color=C_WHITE, align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(1.2), Inches(3.9), Inches(10.9), Inches(0.6),
              "AI Document Analysis · file-processor",
              font_size=14, color=RGBColor(0xCC, 0xBB, 0xFF), align=PP_ALIGN.CENTER)
    ts = datetime.datetime.now().strftime("%Y-%m-%d")
    _text_box(slide, Inches(1.2), Inches(4.55), Inches(10.9), Inches(0.45),
              ts, font_size=12, color=RGBColor(0xAA, 0x99, 0xDD), align=PP_ALIGN.CENTER)


def build_pptx(info: dict, source_filename: str = "document") -> bytes:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_cover(prs, info, source_filename)
    _slides_summary(prs, info)
    _slides_topics(prs, info)
    _slides_entities_dates(prs, info)
    _slides_actions(prs, info)
    _slides_extra_sections(prs, info)
    _slide_closing(prs, source_filename)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()